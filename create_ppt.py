from pptx import Presentation

# Create presentation
prs = Presentation()

# Slide 1: Introduction
slide_layout = prs.slide_layouts[1]
slide = prs.slides.add_slide(slide_layout)
title = slide.shapes.title
content = slide.placeholders[1]
title.text = "Introduction to HCLTech"
content.text = ("Company Name: HCL Technologies Limited (HCLTech)\n"
                "Headquarters: Noida, India\n"
                "Founding Year: 1976\n"
                "Global Presence: 60 countries, 223,000+ employees\n"
                "Vision/Mission: HCLFoundation nurtures clean, green, healthy communities, empowering people to reach full potential.")

# Slide 2: Leadership & Business Structure
slide = prs.slides.add_slide(slide_layout)
title = slide.shapes.title
content = slide.placeholders[1]
title.text = "Leadership & Business Structure"
content.text = ("Chairman Emeritus & Advisor: Shiv Nadar\n"
                "Chairperson: Roshni Nadar Malhotra\n"
                "CEO & MD: C Vijayakumar (since 2016, reappointed 2026–2030)\n\n"
                "Business Segments:\n"
                "• IT and Business Services (ITBS)\n"
                "• Engineering & R&D Services (ERS)\n"
                "• Products & Platforms (P&P)")

# Slide 3: Key Areas & Technologies
slide = prs.slides.add_slide(slide_layout)
title = slide.shapes.title
content = slide.placeholders[1]
title.text = "Key Areas & Technologies"
content.text = ("• Digital Engineering: Cloud, data, AI, 5G\n"
                "• AI & Automation: AI Force, AI Labs, ML, CV, RPA\n"
                "• Cloud Computing: Cloud-native, adoption & migration\n"
                "• Cybersecurity: Risk mitigation & resilience\n"
                "• Other Tech: Blockchain, Analytics, NLP, LLM, MLOps")

# Slide 4: Strengths & Achievements
slide = prs.slides.add_slide(slide_layout)
title = slide.shapes.title
content = slide.placeholders[1]
title.text = "Strengths & Achievements"
content.text = ("• 3rd largest India-headquartered IT services company\n"
                "• Strong revenue & profit growth\n"
                "• Multiple awards & recognitions globally\n"
                "• Innovation-driven culture\n"
                "• ESG & Sustainability leader through HCLFoundation")

# Slide 5: CSR & Future Outlook
slide = prs.slides.add_slide(slide_layout)
title = slide.shapes.title
content = slide.placeholders[1]
title.text = "CSR & Future Outlook"
content.text = ("• HCLFoundation: CSR arm in India\n"
                "• Focus Areas: Environment, education, health, disaster relief\n"
                "• Long-Term Vision: Leverage emerging tech, expand into high-growth sectors, hybrid AI+human execution model.")

# Save Presentation
prs.save("HCLTech_Interview_PPT.pptx")
